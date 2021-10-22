# pip install python-pptx
#pip install Wand

# to watermark the nike_black.png

#magick composite -gravity NorthWest nike_black.png image1.jpg output1.jpg


import os
from pptx import Presentation
from pptx.util import Inches


from pptx import Presentation
from pptx.util import Inches
import os

#list1 = ['output1.jpg','output2.jpg','output3.jpg','output4.jpg','output5.jpg']



prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[8])
placeholder = slide.placeholders[1]
title=slide.shapes.title # assigning a title
subtitle=slide.placeholders[1]
title.text="Sample Title 1"
subtitle.text="Sample Subtitle 1"

os.chdir(r"G:\Assignment\Lib\site-packages\setuptools-58.2.0.dist-info\images")
picture = placeholder.insert_picture('output1.jpg')



prs.save("myppt.pptx")
os.startfile("myppt.pptx")








lyt=prs.slide_layouts[0] # choosing a slide layout
slide=prs.slides.add_slide(lyt) # adding a slide
title=slide.shapes.title # assigning a title
subtitle=slide.placeholders[1] # placeholder for subtitle
title.text="Sample Title 1" # title


#image = slide.placeholders[1].insert_picture("output1.jpg")
subtitle.text="Sample Subtitle 1" # subtitle


os.chdir(r"C:\\Users\\Mangesh\\Desktop\\Assigment_")
picture = placeholder.insert_picture('output1.jpg')

prs.save("Assignment.pptx") # saving file
os.startfile("Assignment.pptx")