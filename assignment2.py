# pip install python-pptx
#pip install Wand

# to watermark the nike_black.png

#magick composite -gravity NorthWest nike_black.png image1.jpg output1.jpg


import os
from pptx import Presentation
from pptx.util import Inches


from pptx import Presentation
from pptx.util import Inches
import os





pr1 = Presentation()
slide1 = pr1.slides.add_slide(pr1.slide_layouts[8])
placeholder = slide1.placeholders[1]
title1=slide1.shapes.title # assigning a title
subtitle1=slide1.placeholders[1]
title1.text="Assignment"
subtitle1.text="Assignment for creating ppt using pptx"

#os.chdir(r"G:\Assignment\Lib\site-packages\setuptools-58.2.0.dist-info\images")
#picture = placeholder.insert_picture('output1.jpg')

#Slide for Picture1
slide2 = pr1.slides.add_slide(pr1.slide_layouts[1])
#title2=slide2.shapes.title
#title2.text="Sample Title 1"

shapes = slide2.shapes
body_shape = shapes.placeholders[1]

title_shape = shapes.title
title_shape.text = 'Sample Title 1'

tf = body_shape.text_frame
tf.text = 'Sample Subtitle 1'

img1 = "output1.jpg"

from_left = Inches(1)
from_right = Inches(3)
from_top = Inches(2.6)
add_picture = slide2.shapes.add_picture(img1,from_left,from_top,from_right)


#Slide for Picture2
slide3 = pr1.slides.add_slide(pr1.slide_layouts[1])

shapes = slide3.shapes
body_shape = shapes.placeholders[1]

title_shape = shapes.title
title_shape.text = 'Sample Title 2'

tf = body_shape.text_frame
tf.text = 'Sample Subtitle 2'
img2 = "output2.jpg"
from_left = Inches(1)
from_right = Inches(3)
from_top = Inches(2.6)
add_picture = slide3.shapes.add_picture(img2,from_left,from_top,from_right)


#Slide for Picture3
slide4 = pr1.slides.add_slide(pr1.slide_layouts[1])
#title4=slide4.shapes.title
#title4.text="Sample Title 3"
shapes = slide4.shapes
body_shape = shapes.placeholders[1]

title_shape = shapes.title
title_shape.text = 'Sample Title 3'

tf = body_shape.text_frame
tf.text = 'Sample Subtitle 3'

img3 = "output3.jpg"

from_left = Inches(1)
from_right = Inches(3)
from_top = Inches(2.6)

add_picture = slide4.shapes.add_picture(img3,from_left,from_top,from_right)



#Slide for Picture4
slide5 = pr1.slides.add_slide(pr1.slide_layouts[1])
#title5=slide5.shapes.title
#title5.text="Sample Title 4"
shapes = slide5.shapes
body_shape = shapes.placeholders[1]

title_shape = shapes.title
title_shape.text = 'Sample Title 4'

tf = body_shape.text_frame
tf.text = 'Sample Subtitle 4'


img4 = "output4.jpg"

from_left = Inches(1)
from_right = Inches(3)
from_top = Inches(2.6)

add_picture = slide5.shapes.add_picture(img4,from_left,from_top,from_right)


#Slide for Picture5
slide6 = pr1.slides.add_slide(pr1.slide_layouts[1])
#title6=slide6.shapes.title
#title6.text="Sample Title 5"
shapes = slide6.shapes
body_shape = shapes.placeholders[1]

title_shape = shapes.title
title_shape.text = 'Sample Title 5'

tf = body_shape.text_frame
tf.text = 'Sample Subtitle 5'

img5 = "output5.jpg"

from_left = Inches(1)
from_right = Inches(3)
from_top = Inches(2.6)

add_picture = slide6.shapes.add_picture(img5,from_left,from_top,from_right)

#pr1.save("assignment2.pptx")
pr1.save("project_assignment.pptx")
#os.startfile("assignment2.pptx")


